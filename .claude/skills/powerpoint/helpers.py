"""PowerPoint generator helpers for SURFStack /powerpoint skill.

Usage from a build script:

    import sys
    sys.path.insert(0, "/path/to/.claude/skills/powerpoint")
    from helpers import Deck

    d = Deck(theme="dark", footer="SURFStack — Topic")
    d.title("...", subtitle="...", kicker="...")
    d.section(num="01", tag="1", title="...", sub="...")
    d.bullets("...", ["a","b","c"])
    ...
    d.save("/path/to/output.pptx")
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Themes ────────────────────────────────────────────────────────────
DARK = dict(
    BG=RGBColor(0x0E, 0x11, 0x16),
    SECTION=RGBColor(0x11, 0x16, 0x1E),
    CARD=RGBColor(0x16, 0x1B, 0x23),
    LINE=RGBColor(0x1F, 0x24, 0x2D),
    FG=RGBColor(0xE8, 0xEC, 0xF1),
    MUTED=RGBColor(0x8A, 0x93, 0xA3),
    ACCENT=RGBColor(0x7C, 0xC4, 0xFF),
    ACCENT2=RGBColor(0xFF, 0xB8, 0x6B),
    WARN=RGBColor(0xFF, 0x7A, 0x7A),
    GOOD=RGBColor(0x7B, 0xE3, 0xA4),
    SECTION_NUM=RGBColor(0x1A, 0x22, 0x30),
    QUOTE_BG=RGBColor(0x14, 0x1C, 0x28),
)
LIGHT = dict(
    BG=RGBColor(0xFA, 0xFB, 0xFD),
    SECTION=RGBColor(0xEE, 0xF1, 0xF6),
    CARD=RGBColor(0xFF, 0xFF, 0xFF),
    LINE=RGBColor(0xD8, 0xDD, 0xE5),
    FG=RGBColor(0x16, 0x1B, 0x23),
    MUTED=RGBColor(0x6B, 0x73, 0x82),
    ACCENT=RGBColor(0x1F, 0x6F, 0xEB),
    ACCENT2=RGBColor(0xC0, 0x6A, 0x16),
    WARN=RGBColor(0xC4, 0x3B, 0x3B),
    GOOD=RGBColor(0x1A, 0x7F, 0x4D),
    SECTION_NUM=RGBColor(0xE2, 0xE6, 0xEC),
    QUOTE_BG=RGBColor(0xEE, 0xF4, 0xFD),
)

DEFAULT_FONT = "Yu Gothic"


class Deck:
    def __init__(self, theme="dark", aspect="16:9", font=DEFAULT_FONT, footer=""):
        self.prs = Presentation()
        if aspect == "16:9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        elif aspect == "4:3":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
        else:
            raise ValueError("aspect must be '16:9' or '4:3'")
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        self.layout = self.prs.slide_layouts[6]  # blank
        self.font = font
        self.footer_text = footer
        self.t = DARK if theme == "dark" else LIGHT
        self.theme_name = theme
        self._slides = []  # for footer page numbers
        self._builders = []  # deferred footer rendering

    # ── primitives ──────────────────────────────────────────────────
    def _new(self, bg=None):
        s = self.prs.slides.add_slide(self.layout)
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.SW, self.SH)
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg or self.t["BG"]
        rect.shadow.inherit = False
        self._slides.append(s)
        return s

    def _text(self, slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
        """runs: list of dicts {text, size, bold, color, font} OR single str."""
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        if isinstance(runs, str):
            runs = [{"text": runs}]
        p = tf.paragraphs[0]
        p.alignment = align
        for ropts in runs:
            r = p.add_run()
            r.text = ropts["text"]
            f = r.font
            f.name = ropts.get("font", self.font)
            f.size = Pt(ropts.get("size", 18))
            f.bold = ropts.get("bold", False)
            f.color.rgb = ropts.get("color", self.t["FG"])
        return tb

    def _paras(self, slide, x, y, w, h, paragraphs, *,
               anchor=MSO_ANCHOR.TOP):
        """paragraphs: list of dicts with optional 'runs' or 'text', 'space_before', 'align'."""
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        for i, spec in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = spec.get("align", PP_ALIGN.LEFT)
            if "space_before" in spec:
                p.space_before = Pt(spec["space_before"])
            runs = spec.get("runs") or [{"text": spec.get("text", "")}]
            for ropts in runs:
                r = p.add_run()
                r.text = ropts["text"]
                f = r.font
                f.name = ropts.get("font", self.font)
                f.size = Pt(ropts.get("size", spec.get("size", 18)))
                f.bold = ropts.get("bold", spec.get("bold", False))
                f.color.rgb = ropts.get("color", spec.get("color", self.t["FG"]))
        return tb

    def _kicker(self, slide, txt):
        if not txt:
            return
        self._text(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4),
                   [{"text": txt.upper(), "size": 12, "bold": True,
                     "color": self.t["MUTED"]}])

    def _h2(self, slide, txt, *, color=None, top=1.0):
        self._text(slide, Inches(0.9), Inches(top), Inches(11.5), Inches(0.9),
                   [{"text": txt, "size": 32, "bold": True,
                     "color": color or self.t["ACCENT"]}])

    def _bullets(self, slide, x, y, w, h, items, *, size=20, gap=8):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_top = 0
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if i > 0:
                p.space_before = Pt(gap)
            if isinstance(it, str):
                runs = [{"text": "• " + it}]
            else:
                runs = [{"text": "• "}] + list(it)
            for ropts in runs:
                r = p.add_run()
                r.text = ropts["text"]
                r.font.name = ropts.get("font", self.font)
                r.font.size = Pt(ropts.get("size", size))
                r.font.bold = ropts.get("bold", False)
                r.font.color.rgb = ropts.get("color", self.t["FG"])

    def _card(self, slide, x, y, w, h, title, body, *, title_color=None):
        title_color = title_color or self.t["ACCENT2"]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.adjustments[0] = 0.06
        box.line.color.rgb = self.t["LINE"]
        box.line.width = Pt(0.75)
        box.fill.solid(); box.fill.fore_color.rgb = self.t["CARD"]
        tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                       w - Inches(0.5), h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = self.font; r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = title_color
        body = [body] if isinstance(body, str) else body
        for line in body:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(6)
            rr = p2.add_run(); rr.text = line
            rr.font.name = self.font; rr.font.size = Pt(14)
            rr.font.color.rgb = self.t["FG"]

    def _quote(self, slide, x, y, w, h, text, *, size=22):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(45000), h)
        bar.line.fill.background()
        bar.fill.solid(); bar.fill.fore_color.rgb = self.t["ACCENT"]
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    x + Emu(45000), y, w - Emu(45000), h)
        bg.line.fill.background()
        bg.fill.solid(); bg.fill.fore_color.rgb = self.t["QUOTE_BG"]
        tb = slide.shapes.add_textbox(x + Inches(0.3), y, w - Inches(0.4), h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.name = self.font; r.font.size = Pt(size)
        r.font.color.rgb = self.t["FG"]

    def _table(self, slide, x, y, w, h, rows):
        nr, nc = len(rows), len(rows[0])
        tbl = slide.shapes.add_table(nr, nc, x, y, w, h).table
        for ci in range(nc):
            tbl.columns[ci].width = w // nc
        for ri, row in enumerate(rows):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    self.t["SECTION"] if ri == 0 else self.t["BG"])
                tf = cell.text_frame
                tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
                tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
                tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = ""
                r = p.add_run(); r.text = cell_text
                r.font.name = self.font
                r.font.size = Pt(14)
                r.font.bold = (ri == 0)
                r.font.color.rgb = self.t["ACCENT"] if ri == 0 else self.t["FG"]

    # ── public API ──────────────────────────────────────────────────
    def title(self, title, *, subtitle="", kicker=""):
        s = self._new()
        if kicker:
            self._text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                       [{"text": kicker.upper(), "size": 14, "bold": True,
                         "color": self.t["MUTED"]}])
        # split title by \n manually
        for i, line in enumerate(title.split("\n")):
            self._text(s, Inches(0.9), Inches(3.0 + i * 1.0), Inches(11.5),
                       Inches(1.1),
                       [{"text": line, "size": 56, "bold": True,
                         "color": self.t["FG"]}])
        if subtitle:
            self._text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.5),
                       [{"text": subtitle, "size": 18,
                         "color": self.t["MUTED"]}])
        return s

    def section(self, *, num, tag, title, sub=""):
        s = self._new(self.t["SECTION"])
        self._text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5),
                   [{"text": f"SECTION {tag}", "size": 14, "bold": True,
                     "color": self.t["ACCENT"]}])
        self._text(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6),
                   [{"text": num, "size": 160, "bold": True,
                     "color": self.t["SECTION_NUM"]}])
        self._text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.2),
                   [{"text": title, "size": 44, "bold": True,
                     "color": self.t["FG"]}])
        if sub:
            self._text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.6),
                       [{"text": sub, "size": 16, "color": self.t["MUTED"]}])
        return s

    def bullets(self, heading, items, *, kicker="", size=20):
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        self._bullets(s, Inches(0.9), Inches(2.2 if kicker else 2.5),
                      Inches(11.5), Inches(4.2), items, size=size)
        return s

    def two_col(self, heading, *, left, right, kicker=""):
        """left/right = (title, [items])."""
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        ly = Inches(2.2 if kicker else 2.5)
        # left col
        self._text(s, Inches(0.9), ly, Inches(5.5), Inches(0.6),
                   [{"text": left[0], "size": 24, "bold": True,
                     "color": self.t["ACCENT2"]}])
        self._bullets(s, Inches(0.9), ly + Inches(0.8), Inches(5.5),
                      Inches(3.5), left[1], size=18)
        # right col
        self._text(s, Inches(7.0), ly, Inches(5.5), Inches(0.6),
                   [{"text": right[0], "size": 24, "bold": True,
                     "color": self.t["ACCENT2"]}])
        self._bullets(s, Inches(7.0), ly + Inches(0.8), Inches(5.5),
                      Inches(3.5), right[1], size=18)
        return s

    def cards(self, heading, items, *, kicker="", cols=3):
        """items: list of (title, body) tuples."""
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        gx, gy = Inches(0.9), Inches(2.0 if kicker else 2.3)
        gap = Inches(0.2)
        total_w = Inches(11.5)
        cw = Emu(int((total_w - gap * (cols - 1)) / cols))
        ch = Inches(1.55)
        for i, (t, b) in enumerate(items):
            col, row = i % cols, i // cols
            x = gx + col * (cw + gap)
            y = gy + row * (ch + gap)
            self._card(s, x, y, cw, ch, t, b)
        return s

    def quote(self, heading, text, *, kicker="", caption=""):
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        self._quote(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.0),
                    text, size=22)
        if caption:
            self._text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6),
                       [{"text": caption, "size": 16,
                         "color": self.t["MUTED"]}])
        return s

    def big(self, heading, big_text, *, caption="", kicker=""):
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        self._text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5),
                   [{"text": big_text, "size": 80, "bold": True,
                     "color": self.t["FG"]}])
        if caption:
            self._text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.4),
                       [{"text": caption, "size": 18, "color": self.t["FG"]}])
        return s

    def table(self, heading, rows, *, kicker="", note=""):
        s = self._new()
        self._kicker(s, kicker)
        self._h2(s, heading, top=0.7 if kicker else 1.0)
        self._table(s, Inches(0.9), Inches(2.0 if kicker else 2.2),
                    Inches(11.5), Inches(4.5), rows)
        if note:
            self._text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
                       [{"text": note, "size": 14, "color": self.t["MUTED"]}])
        return s

    def cta(self, heading, steps, *, kicker="CALL TO ACTION", footer_note=""):
        """steps: list of (title, sub) tuples."""
        s = self._new(RGBColor(0x0E, 0x16, 0x20)
                      if self.theme_name == "dark" else self.t["SECTION"])
        self._kicker(s, kicker)
        self._h2(s, heading)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0),
                                  Inches(11.5), Inches(4.5))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (head, sub) in enumerate(steps):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if i > 0:
                p.space_before = Pt(10)
            r = p.add_run(); r.text = f"{i+1}. "
            r.font.name = self.font; r.font.size = Pt(20)
            r.font.bold = True; r.font.color.rgb = self.t["ACCENT"]
            r2 = p.add_run(); r2.text = head
            r2.font.name = self.font; r2.font.size = Pt(20)
            r2.font.bold = True; r2.font.color.rgb = self.t["FG"]
            if sub:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                rr = p2.add_run(); rr.text = "    " + sub
                rr.font.name = self.font; rr.font.size = Pt(15)
                rr.font.color.rgb = self.t["MUTED"]
        if footer_note:
            self._text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
                       [{"text": footer_note, "size": 14,
                         "color": self.t["MUTED"]}])
        return s

    def end(self, headline, *, subtitle=""):
        s = self._new()
        self._text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                   [{"text": "END", "size": 14, "bold": True,
                     "color": self.t["MUTED"]}])
        for i, line in enumerate(headline.split("\n")):
            self._text(s, Inches(0.9), Inches(3.0 + i * 1.0), Inches(11.5),
                       Inches(1.1),
                       [{"text": line, "size": 56, "bold": True,
                         "color": self.t["FG"]}])
        if subtitle:
            self._text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.5),
                       [{"text": subtitle, "size": 16,
                         "color": self.t["MUTED"]}])
        return s

    # ── footer rendering on save ───────────────────────────────────
    def _render_footers(self):
        total = len(self._slides)
        for i, s in enumerate(self._slides, 1):
            # progress bar
            full_w = self.SW - Inches(1.8)
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9),
                                    Inches(0.05), full_w, Emu(20000))
            bg.line.fill.background()
            bg.fill.solid(); bg.fill.fore_color.rgb = self.t["LINE"]
            fg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9),
                                    Inches(0.05),
                                    int(full_w * i / total), Emu(20000))
            fg.line.fill.background()
            fg.fill.solid(); fg.fill.fore_color.rgb = self.t["ACCENT"]
            # footer
            self._text(s, Inches(0.9), Inches(7.05), Inches(8), Inches(0.3),
                       [{"text": self.footer_text, "size": 10,
                         "color": self.t["MUTED"]}])
            self._text(s, Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.3),
                       [{"text": f"{i} / {total}", "size": 10,
                         "color": self.t["MUTED"]}],
                       align=PP_ALIGN.RIGHT)

    def save(self, path):
        self._render_footers()
        self.prs.save(path)
        return path
